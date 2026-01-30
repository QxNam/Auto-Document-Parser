import io

from pptx import Presentation

from adp.configs.logger import worker_logger as logger
from adp.services.parse.base_parse import BaseParse
from adp.services.parse.parse_registry import ParseRegistry


@ParseRegistry.register([".pptx", ".ppt"])
class PowerPointParse(BaseParse):
    """
    PowerPoint Parser using python-pptx
    """

    def __init__(self, table_output: str = "markdown"):
        self.table_output = table_output

    def parse(
        self,
        file_obj: io.BytesIO,
        engine: str = "auto",
        output_format: str = "markdown",
    ) -> str:
        try:
            prs = Presentation(pptx=file_obj)
            full_text = []

            for i, slide in enumerate(prs.slides):
                slide_num = i + 1
                full_text.append(f"\n--- Slide {slide_num} ---\n")

                slide_content = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        if not shape.has_table:
                            slide_content.append(shape.text.strip())

                    if shape.has_table:
                        markdown_table = self._process_table(shape)
                        slide_content.append(markdown_table)

                full_text.append("\n".join(slide_content))
                full_text.append("\n")

            content = "\n".join(full_text)
            return content

        except Exception as e:
            logger.error(f"Failed to parse PowerPoint: {e}")

    def _process_table(self, shape) -> str:
        """
        Helper function: Convert PPT table to Markdown format.
        """
        table = shape.table
        rows = []

        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                cell_text = cell.text_frame.text.strip().replace("\n", " ")
                cells.append(cell_text)

            rows.append(f"| {' | '.join(cells)} |")

            if i == 0:
                cols = len(row.cells)
                separator = f"| {' | '.join(['---'] * cols)} |"
                rows.append(separator)

        return "\n".join(rows)
